import fitz  # PyMuPDF
import json
import re
import argparse
from pathlib import Path
from docx import Document
from pptx import Presentation

# ================= CONFIG =================

SCRIPT_DIR = Path(__file__).resolve().parent
DATA_DIR = (SCRIPT_DIR / "../../data").resolve()
OUTPUT_DIR = SCRIPT_DIR / "output"

CHUNK_SIZE = 500
CHUNK_OVERLAP = 50

PROGRESS_FILE = OUTPUT_DIR / "progress.json"
CHUNKS_FILE = OUTPUT_DIR / "chunks.jsonl"

# =========================================


def parse_args():
    parser = argparse.ArgumentParser(description="Document Chunker")
    parser.add_argument(
        "--redo",
        action="store_true",
        help="Delete existing chunks and progress and re-chunk from scratch"
    )
    return parser.parse_args()


def load_progress():
    if PROGRESS_FILE.exists():
        with open(PROGRESS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_progress(progress):
    with open(PROGRESS_FILE, "w", encoding="utf-8") as f:
        json.dump(progress, f, indent=2)


def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def chunk_text(text: str):
    words = text.split()
    chunks = []
    start = 0

    while start < len(words):
        end = start + CHUNK_SIZE
        chunks.append(" ".join(words[start:end]))
        start += CHUNK_SIZE - CHUNK_OVERLAP

    return chunks


# ---------- TEXT EXTRACTORS ----------

def extract_docx_text(path: Path) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pptx_text(path: Path) -> str:
    prs = Presentation(path)
    texts = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)

    return "\n".join(texts)


# ---------- PROCESSORS ----------

def process_pdf(pdf_path: Path, folder_meta: dict, progress: dict):
    doc_key = str(pdf_path.resolve())
    last_page_done = progress.get(doc_key, -1)

    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"❌ Failed to open PDF: {pdf_path}\n    {e}")
        return

    total_pages = doc.page_count
    print(f"\n📘 PDF: {pdf_path.relative_to(DATA_DIR)} ({total_pages} pages)")

    if last_page_done >= 0:
        print(f"↪ Resuming from page {last_page_done + 1}")

    with open(CHUNKS_FILE, "a", encoding="utf-8") as out:
        for page_num in range(last_page_done + 1, total_pages):
            try:
                page = doc.load_page(page_num)
                text = page.get_text()

                if not text or len(text.strip()) < 30:
                    print(
                        f"⚠️  No extractable text on "
                        f"{pdf_path.name} page {page_num + 1} "
                        f"(OCR needed)"
                    )
                    progress[doc_key] = page_num
                    save_progress(progress)
                    continue

                chunks = chunk_text(clean_text(text))

                for chunk in chunks:
                    record = {
                        "text": chunk,
                        "doc_name": pdf_path.name,
                        "page": page_num + 1,
                        "source_path": str(pdf_path.resolve()),
                        **folder_meta
                    }
                    out.write(json.dumps(record, ensure_ascii=False) + "\n")

                progress[doc_key] = page_num
                save_progress(progress)

            except Exception as e:
                print(
                    f"❌ ERROR at {pdf_path.name} page {page_num + 1}\n"
                    f"    {e}\n"
                    f"➡ Progress saved. Re-run to resume."
                )
                save_progress(progress)
                return

    print(f"✅ Finished PDF: {pdf_path.name}")


def process_docx(docx_path: Path, folder_meta: dict):
    print(f"\n📄 DOCX: {docx_path.relative_to(DATA_DIR)}")

    try:
        text = extract_docx_text(docx_path)
    except Exception as e:
        print(f"❌ Failed to read DOCX: {docx_path}\n    {e}")
        return

    if not text or len(text.strip()) < 30:
        print("⚠️ No usable text found")
        return

    chunks = chunk_text(clean_text(text))

    with open(CHUNKS_FILE, "a", encoding="utf-8") as out:
        for chunk in chunks:
            record = {
                "text": chunk,
                "doc_name": docx_path.name,
                "page": None,
                "source_path": str(docx_path.resolve()),
                **folder_meta
            }
            out.write(json.dumps(record, ensure_ascii=False) + "\n")

    print(f"✅ Finished DOCX: {docx_path.name}")


def process_pptx(pptx_path: Path, folder_meta: dict):
    print(f"\n📊 PPTX: {pptx_path.relative_to(DATA_DIR)}")

    try:
        text = extract_pptx_text(pptx_path)
    except Exception as e:
        print(f"❌ Failed to read PPTX: {pptx_path}\n    {e}")
        return

    if not text or len(text.strip()) < 30:
        print("⚠️ No usable text found")
        return

    chunks = chunk_text(clean_text(text))

    with open(CHUNKS_FILE, "a", encoding="utf-8") as out:
        for chunk in chunks:
            record = {
                "text": chunk,
                "doc_name": pptx_path.name,
                "page": None,
                "source_path": str(pptx_path.resolve()),
                **folder_meta
            }
            out.write(json.dumps(record, ensure_ascii=False) + "\n")

    print(f"✅ Finished PPTX: {pptx_path.name}")


def process_file(path: Path, folder_meta: dict, progress: dict):
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        process_pdf(path, folder_meta, progress)
    elif suffix == ".docx":
        process_docx(path, folder_meta)
    elif suffix == ".pptx":
        process_pptx(path, folder_meta)
    else:
        print(f"⚠️ Unsupported file type: {path}")


# ---------- MAIN ----------

def main():
    args = parse_args()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    if args.redo:
        print("🔁 REDO MODE: clearing previous output")
        if PROGRESS_FILE.exists():
            PROGRESS_FILE.unlink()
        if CHUNKS_FILE.exists():
            CHUNKS_FILE.unlink()
        progress = {}
    else:
        print("▶ RESUME MODE (default)")
        progress = load_progress()

    print(f"\n📂 Scanning data directory: {DATA_DIR}")

    for year_dir in DATA_DIR.iterdir():
        if not year_dir.is_dir():
            continue

        for sem_dir in year_dir.iterdir():
            if not sem_dir.is_dir():
                continue

            for subject_dir in sem_dir.iterdir():
                if not subject_dir.is_dir():
                    continue

                meta_file = subject_dir / "metadata.json"
                if meta_file.exists():
                    with open(meta_file, "r", encoding="utf-8") as f:
                        folder_meta = json.load(f)
                else:
                    print(f"⚠️ Missing metadata.json in {subject_dir}")
                    folder_meta = {}

                subject_name = folder_meta.get("subject", subject_dir.name)
                print(f"\n📂 Subject: {subject_name}")

                files = []
                files.extend(subject_dir.rglob("*.pdf"))
                files.extend(subject_dir.rglob("*.docx"))
                files.extend(subject_dir.rglob("*.pptx"))

                if not files:
                    print("   (no supported files found)")
                    continue

                for file_path in files:
                    process_file(file_path, folder_meta, progress)

    print("\n🎉 Chunking complete")


if __name__ == "__main__":
    main()